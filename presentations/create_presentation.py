#!/usr/bin/env python3
"""
Neuratex AI Präsentation - PowerPoint Generator v2
CI-Farben: Navy #0A1931, Gold #FFC947, Blau #185ADB
Verbesserungen: Logo oben rechts, Fußzeile, mehr visuelle Abwechslung
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# CI Farben
NAVY = RGBColor(0x0A, 0x19, 0x31)
GOLD = RGBColor(0xFF, 0xC9, 0x47)
BLUE = RGBColor(0x18, 0x5A, 0xDB)
TEXT_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
HEADLINE_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BG = RGBColor(0x11, 0x22, 0x40)
DARK_ACCENT = RGBColor(0x1E, 0x29, 0x3B)

# Pfade
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH = os.path.join(SCRIPT_DIR, "../public/assets/images/neuratex-logo.jpg")
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "Neuratex_AI_Praesentation.pptx")

# Slide dimensions
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_background(slide, color):
    """Setzt den Hintergrund einer Folie"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_logo_top_right(slide, logo_path, width=Inches(1.3)):
    """Fügt Logo oben rechts hinzu"""
    if os.path.exists(logo_path):
        slide.shapes.add_picture(
            logo_path,
            SLIDE_WIDTH - width - Inches(0.4),
            Inches(0.3),
            width=width
        )


def add_footer(slide, text="Neuratex AI"):
    """Fügt Fußzeile unten hinzu"""
    # Fußzeilen-Linie
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.1),
        SLIDE_WIDTH, Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # Fußzeilen-Text
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(9), Inches(0.3))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_LIGHT
    p.font.name = "Arial"


def add_section_header(slide, text, left=Inches(0.5), top=Inches(0.5)):
    """Fügt Abschnittsüberschrift in Gold hinzu"""
    txBox = slide.shapes.add_textbox(left, top, Inches(8), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = "Arial"
    p.font.letter_spacing = Pt(2)
    return txBox


def add_title(slide, text, left=Inches(0.5), top=Inches(0.9), font_size=36, color=HEADLINE_WHITE):
    """Fügt Haupttitel hinzu"""
    txBox = slide.shapes.add_textbox(left, top, Inches(8), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Arial"
    return txBox


def add_text(slide, text, left, top, width, height, font_size=16, color=TEXT_LIGHT, bold=False, align=PP_ALIGN.LEFT):
    """Fügt Text hinzu"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    p.alignment = align
    p.line_spacing = 1.4
    return txBox


def add_bullet_list(slide, items, left, top, width, height, font_size=15, color=TEXT_LIGHT):
    """Fügt Bullet-Liste hinzu"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.line_spacing = 1.6
        p.space_after = Pt(6)
    return txBox


def add_card(slide, title, items, left, top, width, height, title_color=GOLD):
    """Fügt eine Karte mit Rahmen hinzu"""
    # Karten-Hintergrund
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = DARK_ACCENT
    card.line.width = Pt(1)

    # Titel
    add_text(slide, title, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
             font_size=14, color=title_color, bold=True)

    # Items
    if items:
        add_bullet_list(slide, items, left + Inches(0.2), top + Inches(0.5),
                       width - Inches(0.4), height - Inches(0.6), font_size=13)

    return card


def add_highlight_box(slide, text, left, top, width, height, font_size=17):
    """Fügt Gold-Highlight-Box hinzu"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

    # Adjust corner radius
    shape.adjustments[0] = 0.1

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Vertikale Zentrierung
    tf.anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = "Arial"
    p.line_spacing = 1.3

    return shape


def add_icon_circle(slide, number, left, top, size=Inches(0.55)):
    """Fügt nummerierte Kreise hinzu"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()

    tf = circle.text_frame
    tf.anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    return circle


def add_vertical_accent_line(slide, left, top, height):
    """Fügt vertikale Akzentlinie hinzu"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(4), height)
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()
    return line


def add_horizontal_divider(slide, left, top, width):
    """Fügt horizontale Trennlinie hinzu"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = DARK_ACCENT
    line.line.fill.background()
    return line


def add_notes(slide, text):
    """Fügt Sprechernotizen hinzu"""
    if text:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = text


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ========== FOLIE 1: TITEL ==========
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, NAVY)

    # Dekorative Elemente
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(0.15), Inches(1.8))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = GOLD
    accent_bar.line.fill.background()

    # Taglines
    add_text(slide1, "Wissen nutzbar machen.", Inches(0.8), Inches(1.8), Inches(8), Inches(0.5),
             font_size=26, color=TEXT_LIGHT)
    add_text(slide1, "Suche eliminieren.", Inches(0.8), Inches(2.35), Inches(8), Inches(0.5),
             font_size=26, color=TEXT_LIGHT)
    add_text(slide1, "Service beschleunigen.", Inches(0.8), Inches(2.9), Inches(8), Inches(0.5),
             font_size=26, color=TEXT_LIGHT)

    # Haupttitel
    add_text(slide1, "Neuratex AI", Inches(0.8), Inches(3.8), Inches(8), Inches(1),
             font_size=52, color=GOLD, bold=True)

    # Autor
    add_text(slide1, "Philipp Koch", Inches(0.8), Inches(5), Inches(8), Inches(0.5),
             font_size=18, color=TEXT_LIGHT)

    # Logo groß
    if os.path.exists(LOGO_PATH):
        slide1.shapes.add_picture(LOGO_PATH, Inches(6.8), Inches(1.5), width=Inches(2.8))

    add_footer(slide1)
    add_notes(slide1, "Ich helfe technischen Unternehmen dabei, Informationen schnell auffindbar zu machen, ohne ihre Systeme umzubauen.")

    # ========== FOLIE 2: DIE REALITÄT ==========
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, NAVY)

    add_section_header(slide2, "Die Realität in technischen Unternehmen")
    add_title(slide2, "Warum es heute oft unnötig langsam ist")

    # Zwei Karten nebeneinander
    add_card(slide2, "Kritisches Wissen liegt verteilt in:",
             ["PDFs und Scans", "Tabellen und Listen", "Zeichnungen und Katalogen",
              "E-Mails, Ordnern, SharePoint", "Köpfen einzelner Mitarbeiter"],
             Inches(0.5), Inches(2.1), Inches(4.3), Inches(2.5))

    add_card(slide2, "Das Ergebnis im Alltag:",
             ["Suchen statt arbeiten", "Rückfragen statt Lösungen",
              "Abhängigkeit von Schlüsselpersonen", "Lange Antwortzeiten im Service"],
             Inches(5.2), Inches(2.1), Inches(4.3), Inches(2.5))

    add_highlight_box(slide2, "Wenn Information existiert, aber nicht auffindbar ist,\nist sie wirtschaftlich wertlos.",
                      Inches(0.5), Inches(5.0), Inches(9), Inches(1.1))

    add_logo_top_right(slide2, LOGO_PATH)
    add_footer(slide2)

    # ========== FOLIE 3: WAS ES KOSTET ==========
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, NAVY)

    add_section_header(slide3, "Was es wirklich kostet")
    add_title(slide3, "Zeitverlust ist nur die Oberfläche")

    # Vertikale Akzentlinie links
    add_vertical_accent_line(slide3, Inches(0.5), Inches(2.2), Inches(2.8))

    # Direkte Kosten
    add_text(slide3, "Direkte Kosten", Inches(0.75), Inches(2.2), Inches(4), Inches(0.4),
             font_size=16, color=GOLD, bold=True)
    add_bullet_list(slide3, [
        "Suchzeit im Service und Vertrieb",
        "Unterbrechungen und Kontextwechsel",
        "Doppelte Arbeit durch wiederholte Recherche"
    ], Inches(0.75), Inches(2.6), Inches(4), Inches(2))

    # Vertikale Akzentlinie rechts
    add_vertical_accent_line(slide3, Inches(5.2), Inches(2.2), Inches(2.8))

    # Indirekte Kosten
    add_text(slide3, "Indirekte Kosten", Inches(5.45), Inches(2.2), Inches(4), Inches(0.4),
             font_size=16, color=GOLD, bold=True)
    add_bullet_list(slide3, [
        "Langsamer Service wirkt inkompetent",
        "Kunden springen ab",
        "Experten werden zum Flaschenhals",
        "Einarbeitung dauert zu lange"
    ], Inches(5.45), Inches(2.6), Inches(4.2), Inches(2.2))

    add_highlight_box(slide3, "Wachstum scheitert selten an Nachfrage.\nEs scheitert an Zugriff auf Wissen.",
                      Inches(0.5), Inches(5.3), Inches(9), Inches(1.1))

    add_logo_top_right(slide3, LOGO_PATH)
    add_footer(slide3)

    # ========== FOLIE 4: KLASSISCHE ANSÄTZE ==========
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, NAVY)

    add_section_header(slide4, "Warum klassische Ansätze nicht reichen")
    add_title(slide4, "Warum Teams, SharePoint und Ablage\ndas Problem nicht lösen", font_size=30)

    # Große Aussage mit Akzent
    add_text(slide4, "Ablage ist nicht Zugriff", Inches(0.5), Inches(2.3), Inches(9), Inches(0.5),
             font_size=22, color=GOLD, bold=True)

    # Horizontale Linie
    add_horizontal_divider(slide4, Inches(0.5), Inches(2.85), Inches(9))

    add_bullet_list(slide4, [
        "Wer nicht weiß, wonach er suchen muss, findet nichts",
        "Gescannte Tabellen und Zeichnungen sind nicht zuverlässig durchsuchbar",
        "Wissen aus gelösten Fällen wird selten standardisiert festgehalten",
        "Ergebnisse landen in Chats, Telefonaten oder WhatsApp – und verschwinden"
    ], Inches(0.5), Inches(3.1), Inches(9), Inches(2), font_size=16)

    add_highlight_box(slide4, "Viele Unternehmen haben Speicher.\nWas fehlt, ist ein Gedächtnis.",
                      Inches(0.5), Inches(5.3), Inches(9), Inches(1.1), font_size=20)

    add_logo_top_right(slide4, LOGO_PATH)
    add_footer(slide4)

    # ========== FOLIE 5: WAS WIR LIEFERN ==========
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, NAVY)

    add_section_header(slide5, "Was wir konkret liefern")
    add_title(slide5, "Keine Tool-Show. Konkrete Ergebnisse.")

    add_text(slide5, "Wir liefern drei Dinge, je nach Ausgangslage:",
             Inches(0.5), Inches(2.0), Inches(9), Inches(0.4), font_size=16, color=TEXT_LIGHT)

    # Drei Schritte mit nummerierten Kreisen
    steps = [
        "Dokumente werden nutzbar gemacht",
        "Erfahrungswissen wird strukturiert gesichert",
        "Zugriff wird drastisch beschleunigt"
    ]

    y_pos = 2.6
    for i, step in enumerate(steps, 1):
        add_icon_circle(slide5, i, Inches(0.5), Inches(y_pos))
        add_text(slide5, step, Inches(1.2), Inches(y_pos + 0.1), Inches(8), Inches(0.5),
                font_size=20, color=HEADLINE_WHITE)
        y_pos += 0.75

    add_highlight_box(slide5, "Wir ersetzen keine Systeme.\nWir bauen eine Zugriffsebene auf Ihren Bestand.",
                      Inches(0.5), Inches(5.3), Inches(9), Inches(1.1))

    add_logo_top_right(slide5, LOGO_PATH)
    add_footer(slide5)

    # ========== FOLIE 6: LÖSUNG 1 ==========
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, NAVY)

    # Nummer prominent
    big_num = slide6.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(1), Inches(1))
    tf = big_num.text_frame
    p = tf.paragraphs[0]
    p.text = "01"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = GOLD

    add_section_header(slide6, "Lösung 1: Altunterlagen nutzbar machen", left=Inches(1.5), top=Inches(0.5))
    add_title(slide6, "Von Scan-Chaos zu Arbeitsdaten", left=Inches(0.5), top=Inches(1.1))

    # Zwei Karten
    add_card(slide6, "Typische Inputs",
             ["Gescannte Kataloge", "PDF-Tabellen", "Stücklisten",
              "Kennlinien", "Ersatzteilseiten", "Zeichnungen mit Positionsnummern"],
             Inches(0.5), Inches(2.0), Inches(4.3), Inches(2.6))

    add_card(slide6, "Outputs, die im Alltag zählen",
             ["Durchsuchbare Dokumente", "Saubere Tabellen (Excel/CSV)",
              "Importdaten für ERP oder Shop", "Strukturierte Daten mit klaren Feldern"],
             Inches(5.2), Inches(2.0), Inches(4.3), Inches(2.6))

    add_highlight_box(slide6, "Jemand findet in 30 Sekunden, was er heute in 15 Minuten sucht.",
                      Inches(0.5), Inches(5.0), Inches(9), Inches(0.95), font_size=18)

    add_logo_top_right(slide6, LOGO_PATH)
    add_footer(slide6)

    # ========== FOLIE 7: LÖSUNG 2 ==========
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, NAVY)

    # Nummer prominent
    big_num = slide7.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(1), Inches(1))
    tf = big_num.text_frame
    p = tf.paragraphs[0]
    p.text = "02"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = GOLD

    add_section_header(slide7, "Lösung 2: Erfahrungswissen sichern", left=Inches(1.5), top=Inches(0.5))
    add_title(slide7, "Wissen geht sonst in Rente", left=Inches(0.5), top=Inches(1.1))

    # Problem-Karte
    add_card(slide7, "Das Problem",
             ["Trefferquote hängt an Personen, nicht am System",
              "Lösungen sind implizit und werden nie dokumentiert",
              "Neue Mitarbeiter müssen fragen statt handeln"],
             Inches(0.5), Inches(2.0), Inches(4.3), Inches(2.0))

    # Lösung-Karte
    add_card(slide7, "Unsere Lösung: Schlanke Fallkarten",
             ["Maschine / Baureihe", "Symptom / Anfrage", "Lösungsschritte",
              "Teile und Hinweise", "Quelle / Referenz"],
             Inches(5.2), Inches(2.0), Inches(4.3), Inches(2.0))

    # Ergebnis-Zeile
    add_horizontal_divider(slide7, Inches(0.5), Inches(4.2), Inches(9))
    add_text(slide7, "Ergebnis:", Inches(0.5), Inches(4.35), Inches(1.2), Inches(0.4),
             font_size=14, color=GOLD, bold=True)
    add_text(slide7, "Wissen bleibt im Unternehmen  •  Einarbeitung wird schneller  •  Qualität wird stabiler",
             Inches(1.7), Inches(4.35), Inches(7.5), Inches(0.4), font_size=14, color=TEXT_LIGHT)

    add_highlight_box(slide7, "Gute Leute sollen Probleme lösen. Nicht Wissen wiederfinden.",
                      Inches(0.5), Inches(5.0), Inches(9), Inches(0.95), font_size=19)

    add_logo_top_right(slide7, LOGO_PATH)
    add_footer(slide7)

    # ========== FOLIE 8: LÖSUNG 3 ==========
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, NAVY)

    # Nummer prominent
    big_num = slide8.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(1), Inches(1))
    tf = big_num.text_frame
    p = tf.paragraphs[0]
    p.text = "03"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = GOLD

    add_section_header(slide8, "Lösung 3: Internes Assistenzsystem", left=Inches(1.5), top=Inches(0.5))
    add_title(slide8, "Schnelle Antworten aus eigenen Quellen", left=Inches(0.5), top=Inches(1.1))

    # Drei Karten in einer Reihe
    add_card(slide8, "Was es ist",
             ["Internes System für Service & Vertrieb",
              "Beantwortet Fragen aus Ihren Dokumenten"],
             Inches(0.5), Inches(2.0), Inches(2.9), Inches(1.8))

    add_card(slide8, "Was es NICHT ist",
             ["Kein Endkundentool",
              "Kein Chatbot ohne Nachweise",
              "Keine Fantasie-Antworten"],
             Inches(3.55), Inches(2.0), Inches(2.9), Inches(1.8))

    add_card(slide8, "Vertrauen durch",
             ["Antworten mit Quellenbezug",
              "Verlinkung zur Originalstelle",
              "Nachvollziehbar für Mitarbeiter"],
             Inches(6.6), Inches(2.0), Inches(2.9), Inches(1.8))

    # Ergebnis-Zeile
    add_horizontal_divider(slide8, Inches(0.5), Inches(4.1), Inches(9))
    add_text(slide8, "Ergebnis:", Inches(0.5), Inches(4.25), Inches(1.2), Inches(0.4),
             font_size=14, color=GOLD, bold=True)
    add_text(slide8, "Minuten statt Viertelstunden  •  Weniger Rückfragen  •  Experten entlastet",
             Inches(1.7), Inches(4.25), Inches(7.5), Inches(0.4), font_size=14, color=TEXT_LIGHT)

    add_logo_top_right(slide8, LOGO_PATH)
    add_footer(slide8)

    # ========== FOLIE 9: VERBINDUNG ==========
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, NAVY)

    add_section_header(slide9, "So verbinden sich die drei Lösungen")
    add_title(slide9, "Ein System. Drei Einstiegspunkte.")

    # Drei verbundene Boxen
    boxes = [
        ("Dokumente", "sind die Basis", Inches(0.5)),
        ("Erfahrungswissen", "macht es praxisnah", Inches(3.55)),
        ("Assistenz", "macht Zugriff schnell", Inches(6.6))
    ]

    for title, subtitle, left in boxes:
        box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.1), Inches(2.9), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = GOLD
        box.line.width = Pt(2)

        add_text(slide9, title, left + Inches(0.1), Inches(2.2), Inches(2.7), Inches(0.5),
                font_size=16, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide9, subtitle, left + Inches(0.1), Inches(2.6), Inches(2.7), Inches(0.4),
                font_size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Verbindungspfeile (einfache Linien)
    for x in [Inches(3.4), Inches(6.45)]:
        arrow = slide9.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(2.5), Inches(0.25), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

    # Einstiegspunkte
    add_text(slide9, "Man kann starten mit:", Inches(0.5), Inches(3.5), Inches(9), Inches(0.4),
             font_size=15, color=HEADLINE_WHITE, bold=True)

    add_bullet_list(slide9, [
        "Altunterlagen → wenn Dokumentchaos dominiert",
        "Wissensfällen → wenn Experten Engpass sind",
        "Assistenz → wenn Suche im Alltag Zeit frisst"
    ], Inches(0.5), Inches(3.9), Inches(9), Inches(1.4), font_size=15)

    add_highlight_box(slide9, "Das ist kein Bauchladen. Es ist eine klare Stufenlogik.",
                      Inches(0.5), Inches(5.6), Inches(9), Inches(0.85), font_size=18)

    add_logo_top_right(slide9, LOGO_PATH)
    add_footer(slide9)

    # ========== FOLIE 10: WARUM JETZT ==========
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, NAVY)

    add_section_header(slide10, "Warum jetzt sinnvoll")
    add_title(slide10, "Timing ist gerade ideal")

    # Vier Punkte als kleine Karten
    reasons = [
        "Neue ERP-/CRM-Einführungen erhöhen\nSichtbarkeit von Datenproblemen",
        "Ohne saubere Dokumentbasis werden\nneue Systeme langsamer akzeptiert",
        "Wissensverlust durch Renteneintritte\nist planbar, aber oft unvorbereitet",
        "Service-Anfragen steigen,\nwährend Fachkräfte knapp sind"
    ]

    positions = [(0.5, 2.1), (5.0, 2.1), (0.5, 3.6), (5.0, 3.6)]

    for i, (reason, (x, y)) in enumerate(zip(reasons, positions)):
        # Kleine Nummer
        num = slide10.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.35), Inches(0.35))
        num.fill.solid()
        num.fill.fore_color.rgb = GOLD
        num.line.fill.background()
        tf = num.text_frame
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER

        add_text(slide10, reason, Inches(x + 0.5), Inches(y), Inches(4), Inches(0.9),
                font_size=14, color=TEXT_LIGHT)

    add_highlight_box(slide10, "Ein neues System ohne nutzbares Wissen ist wie\nein neues Lager ohne Beschriftung.",
                      Inches(0.5), Inches(5.3), Inches(9), Inches(1.1), font_size=17)

    add_logo_top_right(slide10, LOGO_PATH)
    add_footer(slide10)

    # ========== FOLIE 11: VORGEHEN ==========
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, NAVY)

    add_section_header(slide11, "Vorgehen ohne Großprojekt")
    add_title(slide11, "Pragmatisch. Schrittweise. Messbar.")

    steps = [
        ("1", "Machbarkeitscheck", "An echten Unterlagen und echten Fragen"),
        ("2", "Kleines Paket", "Mit sichtbarem Output"),
        ("3", "Skalierung", "Auf weitere Fabrikate oder Dokumenttypen"),
        ("4", "Integration", "Optional in bestehende Systeme")
    ]

    y_pos = 2.2
    for num, title, desc in steps:
        add_icon_circle(slide11, num, Inches(0.5), Inches(y_pos))
        add_text(slide11, title, Inches(1.2), Inches(y_pos), Inches(3), Inches(0.4),
                font_size=18, color=GOLD, bold=True)
        add_text(slide11, desc, Inches(1.2), Inches(y_pos + 0.35), Inches(8), Inches(0.35),
                font_size=14, color=TEXT_LIGHT)
        y_pos += 0.85

    add_highlight_box(slide11, "Wir beweisen Nutzen zuerst. Danach wird entschieden.",
                      Inches(0.5), Inches(5.7), Inches(9), Inches(0.8), font_size=18)

    add_logo_top_right(slide11, LOGO_PATH)
    add_footer(slide11)

    # ========== FOLIE 12: WAS WIR BRAUCHEN ==========
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12, NAVY)

    add_section_header(slide12, "Was wir von Ihnen brauchen")
    add_title(slide12, "Minimaler Aufwand auf Kundenseite")

    add_text(slide12, "Für den Start reichen:", Inches(0.5), Inches(2.2), Inches(9), Inches(0.4),
             font_size=18, color=GOLD, bold=True)

    items = [
        "2–5 Dokumente, die heute nerven",
        "3–10 typische Fragen aus dem Alltag",
        "Ein Ansprechpartner aus Service oder Produktmanagement"
    ]

    y_pos = 2.8
    for item in items:
        # Checkmark-artige Box
        check = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(0.35), Inches(0.35))
        check.fill.solid()
        check.fill.fore_color.rgb = GOLD
        check.line.fill.background()
        tf = check.text_frame
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = "✓"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER

        add_text(slide12, item, Inches(1.0), Inches(y_pos + 0.05), Inches(8), Inches(0.4),
                font_size=18, color=HEADLINE_WHITE)
        y_pos += 0.65

    # Großes "Das war's"
    add_text(slide12, "Das war's.", Inches(0.5), Inches(5.0), Inches(9), Inches(0.8),
             font_size=42, color=GOLD, bold=True)

    add_logo_top_right(slide12, LOGO_PATH)
    add_footer(slide12)

    # ========== FOLIE 13: ERGEBNISSE ==========
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13, NAVY)

    add_section_header(slide13, "Was Sie nach kurzer Zeit sehen")
    add_title(slide13, "Konkrete Ergebnisse statt Folien")

    add_text(slide13, "Beispielhafte Outcomes:", Inches(0.5), Inches(2.2), Inches(9), Inches(0.4),
             font_size=16, color=GOLD, bold=True)

    outcomes = [
        "Dokumente, die wirklich auffindbar sind",
        "Tabellen, die man direkt weiterverwenden kann",
        "Erste Wissensfälle, die jeder findet",
        "Eine Such- oder Assistenzoberfläche, die im Alltag testbar ist"
    ]

    y_pos = 2.8
    for i, outcome in enumerate(outcomes, 1):
        add_icon_circle(slide13, i, Inches(0.5), Inches(y_pos))
        add_text(slide13, outcome, Inches(1.2), Inches(y_pos + 0.1), Inches(8), Inches(0.5),
                font_size=17, color=HEADLINE_WHITE)
        y_pos += 0.7

    add_logo_top_right(slide13, LOGO_PATH)
    add_footer(slide13)

    # ========== FOLIE 14: SICHERHEIT ==========
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14, NAVY)

    add_section_header(slide14, "Sicherheit und Vertraulichkeit")
    add_title(slide14, "Wenn Daten Betriebsgeheimnis sind,\nwird das mitgedacht", font_size=30)

    # Sicherheitsaspekte als Karten-Grid
    security_items = [
        ("Interne Nutzung", "Keine externe Weitergabe"),
        ("Zugriffsrechte", "Rollenbasierte Kontrolle"),
        ("Transparenz", "Sichtbar welche Quellen genutzt werden"),
        ("Keine Veröffentlichung", "Inhalte bleiben geschützt"),
        ("Klare Abgrenzung", "Was rein darf und was nicht")
    ]

    positions = [(0.5, 2.4), (3.55, 2.4), (6.6, 2.4), (2.0, 3.7), (5.1, 3.7)]

    for (title, desc), (x, y) in zip(security_items, positions):
        card = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.9), Inches(1.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = DARK_ACCENT

        add_text(slide14, title, Inches(x + 0.15), Inches(y + 0.1), Inches(2.6), Inches(0.35),
                font_size=13, color=GOLD, bold=True)
        add_text(slide14, desc, Inches(x + 0.15), Inches(y + 0.45), Inches(2.6), Inches(0.4),
                font_size=11, color=TEXT_LIGHT)

    add_highlight_box(slide14, "Wissen soll verfügbar sein. Aber nur für die Richtigen.",
                      Inches(0.5), Inches(5.1), Inches(9), Inches(0.85), font_size=18)

    add_logo_top_right(slide14, LOGO_PATH)
    add_footer(slide14)

    # ========== FOLIE 15: ABSCHLUSS ==========
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15, NAVY)

    # Dekorative Elemente
    accent_bar = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(0.15), Inches(2))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = GOLD
    accent_bar.line.fill.background()

    add_text(slide15, "Wir ersetzen keine Menschen.", Inches(0.8), Inches(1.8), Inches(8), Inches(0.6),
             font_size=24, color=TEXT_LIGHT)

    add_text(slide15, "Wir ersetzen Sucharbeit,\nUnterbrechungen und\nWissensverlust.",
             Inches(0.8), Inches(2.5), Inches(8), Inches(1.5),
             font_size=32, color=GOLD, bold=True)

    # Kontakt-Box
    contact_box = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.0), Inches(4), Inches(1.5))
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = CARD_BG
    contact_box.line.color.rgb = DARK_ACCENT

    add_text(slide15, "Kontakt", Inches(0.7), Inches(5.1), Inches(3.5), Inches(0.3),
             font_size=12, color=GOLD, bold=True)
    add_text(slide15, "Philipp Koch", Inches(0.7), Inches(5.4), Inches(3.5), Inches(0.35),
             font_size=18, color=HEADLINE_WHITE, bold=True)
    add_text(slide15, "Neuratex AI", Inches(0.7), Inches(5.75), Inches(3.5), Inches(0.3),
             font_size=14, color=TEXT_LIGHT)
    add_text(slide15, "philippkoch@neuratex.de", Inches(0.7), Inches(6.05), Inches(3.5), Inches(0.3),
             font_size=13, color=TEXT_LIGHT)

    # Logo groß
    if os.path.exists(LOGO_PATH):
        slide15.shapes.add_picture(LOGO_PATH, Inches(6), Inches(4.3), width=Inches(3.5))

    add_footer(slide15)

    # Speichern
    prs.save(OUTPUT_PATH)
    print(f"Präsentation gespeichert: {OUTPUT_PATH}")


if __name__ == "__main__":
    create_presentation()
